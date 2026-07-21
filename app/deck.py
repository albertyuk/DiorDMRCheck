"""KOL Efficiency Report — .pptx generation (native, editable charts).

One 16:9 slide (13.33"×7.5"), white background, three panels with black
header bars over #F2F2F2. Left column: donut + price bars; right column:
CPM + CPE bars, insights, basis footnote. All charts are native OOXML chart
parts (the team edits decks), built with python-pptx and post-processed at
the chart-XML level for the pieces python-pptx doesn't expose: donut hole
size, per-point label colors (white on dark slices), and sliver labels
pulled outside with leader lines.

After building, ``assert_chart_cache`` re-opens the package and diffs every
cached chart value against the computed metrics — a deck whose XML disagrees
with the numbers must never ship.
"""
from __future__ import annotations

import io
import re
import zipfile
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .effreport import COOPS, TEXTS, TIERS

# palette (spec §6)
SOFT_COLOR = RGBColor(0x00, 0x00, 0x00)
PAID_COLOR = RGBColor(0xC0, 0x00, 0x00)
PANEL_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DONUT_COLORS = {  # clockwise from 12 o'clock; must cover every TIER × COOP —
    # the reference deck's file happens to lack TOP PAID, but uploads can
    # contain it, and a group missing here would vanish from the donut.
    "TOP PAID": "8FAADC", "TOP SOFT": "4472C4",
    "MID PAID": "ED7D31", "MID SOFT": "F8CBAD",
    "BOT PAID": "A6A6A6", "BOT SOFT": "D9D9D9",
    "KOC PAID": "203864", "KOC SOFT": "FFC000",
}
UNCLASSIFIED_COLOR = "BFBFBF"
DONUT_ORDER = ["TOP PAID", "TOP SOFT", "MID PAID", "MID SOFT", "BOT PAID",
               "BOT SOFT", "KOC PAID", "KOC SOFT"]


def _luminance(hex6: str) -> float:
    r, g, b = (int(hex6[i:i + 2], 16) / 255 for i in (0, 2, 4))
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


SLIVER_PCT = 3.0
DARK_LUMA = 0.4

C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NSMAP = {"c": C_NS, "a": A_NS}


def _c(tag: str) -> str:
    return f"{{{C_NS}}}{tag}"


# ------------------------------------------------------------ slide helpers

def _panel(slide, x, y, w, h, title: str):
    """Black header bar (white bold caps) over a light-gray panel body."""
    bar_h = Inches(0.32)
    bar = slide.shapes.add_shape(1, x, y, w, bar_h)  # 1 = rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title.upper()
    run.font.bold = True
    run.font.size = Pt(11)
    run.font.color.rgb = WHITE

    body = slide.shapes.add_shape(1, x, y + bar_h, w, h - bar_h)
    body.fill.solid()
    body.fill.fore_color.rgb = PANEL_GRAY
    body.line.fill.background()
    return body


def _text(slide, x, y, w, h, lines: list[str], size=8.5, bold=False,
          color=BLACK, align=PP_ALIGN.LEFT, bullet=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ("• " if bullet else "") + line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Arial"
    return box


def _legend_swatches(slide, x, y):
    for i, (label, color) in enumerate((("SOFT", SOFT_COLOR),
                                        ("PAID", PAID_COLOR))):
        sq = slide.shapes.add_shape(1, x + Inches(i * 0.75), y,
                                    Inches(0.10), Inches(0.10))
        sq.fill.solid()
        sq.fill.fore_color.rgb = color
        sq.line.fill.background()
        _text(slide, x + Inches(i * 0.75 + 0.12), y - Inches(0.025),
              Inches(0.6), Inches(0.16), [label], size=7.5, bold=True)


def _style_bar_chart(chart, number_format: str):
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 120
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = number_format
    dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(8)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    cat = chart.category_axis
    cat.has_major_gridlines = False
    cat.tick_labels.font.size = Pt(9)
    cat.tick_labels.font.bold = True
    cat.tick_labels.font.color.rgb = BLACK
    cat.format.line.fill.background()
    val = chart.value_axis
    val.visible = False
    val.has_major_gridlines = False
    for series, color in zip(chart.series, (SOFT_COLOR, PAID_COLOR)):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color


def _bar_chart(slide, x, y, w, h, groups: dict, value_key: str,
               number_format: str):
    """SOFT/PAID clustered columns over TOP/MID/BOT/KOC. Missing groups are
    None → a gap in the chart, never a zero bar."""
    data = CategoryChartData()
    data.categories = TIERS
    for coop in ("SOFT", "PAID"):
        data.add_series(coop, tuple(
            (groups.get(f"{t} {coop}") or {}).get(value_key) for t in TIERS))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, data)
    _style_bar_chart(gf.chart, number_format)
    return gf


def _donut_chart(slide, x, y, w, h, slices: list[tuple[str, float]],
                 share_decimals: int):
    data = CategoryChartData()
    data.categories = [s[0] for s in slices]
    data.add_series("share", tuple(s[1] for s in slices))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, w, h, data)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(7.5)
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    fmt = "0." + "0" * share_decimals + '"%"'
    dl.number_format = fmt
    dl.number_format_is_linked = False
    dl.font.size = Pt(8)
    dl.font.bold = True
    series = chart.series[0]
    for i, (label, _val) in enumerate(slices):
        pt = series.points[i]
        pt.format.fill.solid()
        hex6 = DONUT_COLORS.get(label, UNCLASSIFIED_COLOR)
        pt.format.fill.fore_color.rgb = RGBColor.from_string(hex6)
    return gf


# --------------------------------------------------- chart XML post-passes

def _lbl_txpr(parent, white: bool = False):
    """a:txPr with an 8pt bold defRPr (optionally white) — label text props."""
    txpr = etree.SubElement(parent, _c("txPr"))
    etree.SubElement(txpr, f"{{{A_NS}}}bodyPr")
    etree.SubElement(txpr, f"{{{A_NS}}}lstStyle")
    p = etree.SubElement(txpr, f"{{{A_NS}}}p")
    ppr = etree.SubElement(p, f"{{{A_NS}}}pPr")
    rpr = etree.SubElement(ppr, f"{{{A_NS}}}defRPr")
    rpr.set("sz", "800")
    rpr.set("b", "1")
    if white:
        fill = etree.SubElement(rpr, f"{{{A_NS}}}solidFill")
        clr = etree.SubElement(fill, f"{{{A_NS}}}srgbClr")
        clr.set("val", "FFFFFF")
    etree.SubElement(p, f"{{{A_NS}}}endParaRPr")
    return txpr


def _postprocess_donut_xml(chart_part, slices: list[tuple[str, float]]):
    """python-pptx has no API for these — patch the chart XML directly:
    hole size ≈58%, white labels on dark slices, and sliver labels (<~3%)
    pulled outside the ring with leader lines enabled."""
    root = chart_part._element  # chartSpace lxml element
    ser = root.find(f".//{_c('ser')}")
    donut = root.find(f".//{_c('doughnutChart')}")
    hole = donut.find(_c("holeSize"))
    if hole is None:
        hole = etree.SubElement(donut, _c("holeSize"))
    hole.set("val", "58")

    # Per-point label overrides live in a *series*-level dLbls (python-pptx
    # only writes the plot-level one). Schema position: after dPt, before cat.
    dlbls = ser.find(_c("dLbls"))
    if dlbls is None:
        dlbls = etree.Element(_c("dLbls"))
        anchor = ser.find(_c("cat"))
        if anchor is None:
            anchor = ser.find(_c("val"))
        anchor.addprevious(dlbls)
    n_dlbl = len(dlbls.findall(_c("dLbl")))
    total = sum(v for _, v in slices) or 1

    # share format — reused verbatim in every per-point override below, since
    # PowerPoint does not reliably inherit the ser-level numFmt into a c:dLbl
    # that carries its own layout/txPr (labels render as bare "7.9" without it)
    fmt = '0.0"%"'
    plot_dlbls = donut.find(_c("dLbls"))
    if plot_dlbls is not None:
        nf = plot_dlbls.find(_c("numFmt"))
        if nf is not None and nf.get("formatCode"):
            fmt = nf.get("formatCode")

    for i, (label, val) in enumerate(slices):
        hex6 = DONUT_COLORS.get(label, UNCLASSIFIED_COLOR)
        dark = _luminance(hex6) < DARK_LUMA
        sliver = (val / total * 100) < SLIVER_PCT
        if not (dark or sliver):
            continue
        dlbl = etree.Element(_c("dLbl"))
        dlbls.insert(n_dlbl, dlbl)  # dLbl elements precede series-wide props
        n_dlbl += 1
        idx = etree.SubElement(dlbl, _c("idx"))
        idx.set("val", str(i))
        if sliver:
            # pull the label outside the ring; leader line connects it
            layout = etree.SubElement(dlbl, _c("layout"))
            manual = etree.SubElement(layout, _c("manualLayout"))
            for tag, v in (("x", "0.18"), ("y", "-0.12")):
                el = etree.SubElement(manual, _c(tag))
                el.set("val", v)
        ptfmt = etree.SubElement(dlbl, _c("numFmt"))  # after layout, before txPr
        ptfmt.set("formatCode", fmt)
        ptfmt.set("sourceLinked", "0")
        _lbl_txpr(dlbl, white=dark and not sliver)
        for tag in ("showLegendKey", "showVal", "showCatName", "showSerName",
                    "showPercent", "showBubbleSize"):
            el = etree.SubElement(dlbl, _c(tag))
            el.set("val", "1" if tag == "showVal" else "0")

    # A ser-level dLbls is authoritative once present — renderers fall back to
    # their own defaults (category + series name, theme font) for anything it
    # omits, so it must carry the full series-wide config, not just overrides.
    # Schema order within dLbls: dLbl*, numFmt, txPr, show*, showLeaderLines.
    numfmt = etree.SubElement(dlbls, _c("numFmt"))
    numfmt.set("formatCode", fmt)
    numfmt.set("sourceLinked", "0")
    _lbl_txpr(dlbls)
    for tag, v in (("showLegendKey", "0"), ("showVal", "1"),
                   ("showCatName", "0"), ("showSerName", "0"),
                   ("showPercent", "0"), ("showBubbleSize", "0"),
                   ("showLeaderLines", "1")):
        el = etree.SubElement(dlbls, _c(tag))
        el.set("val", v)


# -------------------------------------------------------------- deck build

def build_deck(analysis: dict) -> bytes:
    cfg = analysis["config"]
    texts = TEXTS.get(cfg["language"], TEXTS["en"])
    groups = analysis["metrics"]["groups"]
    totals = analysis["metrics"]["totals"]
    ins = analysis["insights"]
    basis = cfg["basis"]
    cpm_key = "cpm_pooled" if basis == "pooled" else "cpm_perpost"
    cpe_key = "cpe_pooled" if basis == "pooled" else "cpe_perpost"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    _text(slide, Inches(0.35), Inches(0.12), Inches(9), Inches(0.4),
          [texts["deck_title"]], size=18, bold=True)

    # ---- left column (~34%)
    lx, lw = Inches(0.35), Inches(4.35)
    donut_body = _panel(slide, lx, Inches(0.62), lw, Inches(3.4),
                        texts["donut_title"])
    slices = [(g, groups[g]["share"]) for g in DONUT_ORDER if g in groups]
    if totals["unclassified"]:
        slices.append(("UNCLASSIFIED", totals["unclassified_share"]))
    donut_gf = _donut_chart(
        slide, lx + Inches(0.1), Inches(1.02), lw - Inches(0.2), Inches(2.9),
        slices, cfg["share_decimals"])
    _postprocess_donut_xml(donut_gf.chart.part, slices)

    price_body = _panel(slide, lx, Inches(4.14), lw, Inches(3.0),
                        texts["price_title"])
    _legend_swatches(slide, lx + Inches(0.12), Inches(4.56))
    _bar_chart(slide, lx + Inches(0.05), Inches(4.74), lw - Inches(0.1),
               Inches(1.62), groups, "avg_price", '0.0,"k"')
    _text(slide, lx + Inches(0.12), Inches(6.42), lw - Inches(0.24),
          Inches(0.62), ins["price"], size=7.5, bold=True, bullet=True)

    # ---- right column
    rx, rw = Inches(4.9), Inches(8.08)
    _panel(slide, rx, Inches(0.62), rw, Inches(6.52), texts["eff_title"])
    _legend_swatches(slide, rx + Inches(0.15), Inches(1.06))
    # floor-divide: EMU coordinates are xsd integers — a float here (e.g.
    # cx="3511296.0") makes PowerPoint "repair" the deck by deleting the
    # chart frame, though LibreOffice renders it fine.
    half = (rw - Inches(0.4)) // 2
    _text(slide, rx + Inches(0.15), Inches(1.28), half, Inches(0.2),
          [texts["cpm"]], size=9, bold=True)
    _bar_chart(slide, rx + Inches(0.1), Inches(1.5), half, Inches(2.95),
               groups, cpm_key, "0")
    _text(slide, rx + Inches(0.25) + half, Inches(1.28), half, Inches(0.2),
          [texts["cpe"]], size=9, bold=True)
    _bar_chart(slide, rx + Inches(0.2) + half, Inches(1.5), half, Inches(2.95),
               groups, cpe_key, "0.0")

    bullets = ins["efficiency"] + ins["caveats"]
    _text(slide, rx + Inches(0.15), Inches(4.62), rw - Inches(0.3),
          Inches(1.9), bullets, size=8, bold=True, bullet=True)
    _text(slide, rx + Inches(0.15), Inches(6.62), rw - Inches(0.3),
          Inches(0.44), [ins["footnote"]], size=7, color=RGBColor(0x59, 0x59, 0x59))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ------------------------------------------------------- chart-data assert

def _cached_values(chart_xml: bytes) -> dict[str, list[Optional[float]]]:
    root = etree.fromstring(chart_xml)
    out = {}
    for ser in root.findall(f".//{_c('ser')}"):
        tx = ser.find(f"{_c('tx')}/{_c('strRef')}/{_c('strCache')}/{_c('pt')}/{_c('v')}")
        name = tx.text if tx is not None else f"series{len(out)}"
        vals: dict[int, Optional[float]] = {}
        num_cache = ser.find(f"{_c('val')}/{_c('numRef')}/{_c('numCache')}")
        if num_cache is None:
            continue
        count = int(num_cache.find(_c("ptCount")).get("val"))
        for pt in num_cache.findall(_c("pt")):
            v = pt.find(_c("v"))
            vals[int(pt.get("idx"))] = float(v.text) if v is not None and v.text else None
        out[name] = [vals.get(i) for i in range(count)]
    return out


def assert_chart_cache(pptx_bytes: bytes, analysis: dict) -> None:
    """Parse every embedded chart's cached values out of the finished package
    and diff them against the computed metrics. Raises on any mismatch —
    the deck must not ship if its XML disagrees with the numbers."""
    from .effreport import VerificationError
    cfg = analysis["config"]
    groups = analysis["metrics"]["groups"]
    totals = analysis["metrics"]["totals"]
    cpm_key = "cpm_pooled" if cfg["basis"] == "pooled" else "cpm_perpost"
    cpe_key = "cpe_pooled" if cfg["basis"] == "pooled" else "cpe_perpost"

    def series_for(key):
        return {coop: [(groups.get(f"{t} {coop}") or {}).get(key)
                       for t in TIERS] for coop in COOPS}

    donut_expected = [groups[g]["share"] for g in DONUT_ORDER if g in groups]
    if totals["unclassified"]:
        donut_expected.append(totals["unclassified_share"])
    expected_charts = [
        {"share": donut_expected},
        series_for("avg_price"),
        series_for(cpm_key),
        series_for(cpe_key),
    ]

    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
        chart_names = sorted(
            n for n in z.namelist()
            if re.fullmatch(r"ppt/charts/chart\d+\.xml", n))
        cached = [_cached_values(z.read(n)) for n in chart_names]

    if len(cached) != len(expected_charts):
        raise VerificationError(
            f"expected {len(expected_charts)} charts, found {len(cached)}")
    for want, got, name in zip(expected_charts, cached, chart_names):
        for series, values in want.items():
            gvals = got.get(series)
            if gvals is None:
                raise VerificationError(f"{name}: series {series!r} missing")
            for i, (a, b) in enumerate(zip(values, gvals)):
                if a is None and b is None:
                    continue
                if a is None or b is None or abs(a - b) > 1e-6:
                    raise VerificationError(
                        f"{name} series {series!r} point {i}: "
                        f"computed {a} vs cached {b}")
